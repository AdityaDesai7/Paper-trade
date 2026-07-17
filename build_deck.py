"""Generate a professional 'PetroQuant -> Sparta-class platform' deck."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Design system -------------------------------------------------------
INK      = RGBColor(0x16, 0x23, 0x2E)   # near-black text
PETROL   = RGBColor(0x0C, 0x3B, 0x5B)   # deep petrol navy (primary)
PETROL2  = RGBColor(0x12, 0x4A, 0x6E)   # lighter navy
ACCENT   = RGBColor(0xE9, 0x87, 0x3B)   # oil amber (accent)
TEAL     = RGBColor(0x1F, 0x9E, 0x8F)   # supporting teal
LIGHT    = RGBColor(0xF4, 0xF6, 0xF8)   # light panel bg
CARD     = RGBColor(0xFF, 0xFF, 0xFF)   # white card
MUTED    = RGBColor(0x64, 0x74, 0x83)   # muted gray text
LINE     = RGBColor(0xD9, 0xE0, 0xE6)   # hairline
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

FONT   = "Segoe UI"
FONT_L = "Segoe UI Light"
FONT_SB= "Segoe UI Semibold"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
import os
ENABLE_SHADOW = os.environ.get("DECK_SHADOW", "1") == "1"

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---- helpers -------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shadow=False,
         rounded=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow and ENABLE_SHADOW:
        # schema-safe outer shadow: effectLst must precede scene3d/sp3d/extLst
        spPr = shp._element.spPr
        ef = spPr.makeelement(qn('a:effectLst'), {})
        sh = ef.makeelement(qn('a:outerShdw'),
                            {'blurRad': '80000', 'dist': '38100',
                             'dir': '5400000', 'rotWithShape': '0'})
        clr = sh.makeelement(qn('a:srgbClr'), {'val': '9AA7B2'})
        alpha = clr.makeelement(qn('a:alpha'), {'val': '30000'})
        clr.append(alpha); sh.append(clr); ef.append(sh)
        spPr.insert_element_before(ef, 'a:scene3d', 'a:sp3d', 'a:extLst')
    return shp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text,size,color,bold,font,italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold, fnt, *rest) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = fnt
            if rest and rest[0]:
                r.font.italic = True
    return tb


def bullets(s, x, y, w, h, items, size=15, gap=10, color=INK, marker=ACCENT,
            lead_bold=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0); p.line_spacing = 1.02
        m = p.add_run(); m.text = "\u2013  "
        m.font.size = Pt(size); m.font.color.rgb = marker
        m.font.bold = True; m.font.name = FONT
        # allow "Lead — rest" split
        if "|" in it:
            head, rest = it.split("|", 1)
            r1 = p.add_run(); r1.text = head
            r1.font.size = Pt(size); r1.font.color.rgb = color
            r1.font.bold = True; r1.font.name = FONT
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = MUTED
            r2.font.bold = False; r2.font.name = FONT
        else:
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = False; r.font.name = FONT
    return tb


def header(s, kicker, title, num):
    rect(s, 0, 0, 13.333, 7.5, fill=WHITE)              # base
    rect(s, 0, 0, 0.22, 7.5, fill=PETROL)               # left spine
    rect(s, 0.22, 0, 0.06, 7.5, fill=ACCENT)            # accent line
    txt(s, 0.75, 0.5, 11.8, 0.35,
        [[(kicker.upper(), 12, ACCENT, True, FONT_SB)]], space_after=0)
    txt(s, 0.72, 0.82, 11.8, 0.7,
        [[(title, 27, PETROL, True, FONT)]], space_after=0)
    rect(s, 0.75, 1.52, 1.7, 0.045, fill=ACCENT)
    # footer
    txt(s, 0.75, 7.05, 8, 0.3,
        [[("PetroQuant  ·  Competitive Strategy Brief", 9, MUTED, False, FONT)]],
        space_after=0)
    txt(s, 11.8, 7.05, 0.9, 0.3,
        [[(f"{num:02d}", 9, MUTED, True, FONT)]], align=PP_ALIGN.RIGHT,
        space_after=0)


# =========================================================================
# SLIDE 1 — TITLE
# =========================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=PETROL)
rect(s, 0, 0, 13.333, 7.5, fill=PETROL)
# subtle band
rect(s, 0, 5.55, 13.333, 0.06, fill=ACCENT)
rect(s, 0, 0, 13.333, 0.14, fill=ACCENT)
txt(s, 0.9, 1.15, 11, 0.4,
    [[("ANALYST BRIEF  ·  JULY 2026", 13, ACCENT, True, FONT_SB)]],
    space_after=0)
txt(s, 0.85, 1.9, 11.6, 2.4,
    [[("From PetroQuant to a", 40, WHITE, False, FONT_L)],
     [("Sparta-class Platform", 48, WHITE, True, FONT)]],
    space_after=4, line_spacing=1.0)
txt(s, 0.9, 4.15, 11.2, 1.0,
    [[("What Sparta Commodities is, why it leads oil-trading intelligence, and the",
       16, RGBColor(0xC7,0xD6,0xE2), False, FONT)],
     [("challenges PetroQuant must solve to get there.",
       16, RGBColor(0xC7,0xD6,0xE2), False, FONT)]],
    space_after=2, line_spacing=1.1)
txt(s, 0.9, 5.8, 11, 0.9,
    [[("Scope   ", 12, ACCENT, True, FONT_SB),
      ("Company teardown  ·  USP analysis  ·  gap assessment  ·  roadmap",
       12, RGBColor(0xAE, 0xC0, 0xCE), False, FONT)]],
    space_after=0)

# =========================================================================
# SLIDE 2 — AGENDA
# =========================================================================
s = slide()
header(s, "Contents", "What this brief covers", 2)
ag = [
    ("01", "What Sparta is", "Definition, origin, and category"),
    ("02", "Sparta at a glance", "Facts, funding, clients, footprint"),
    ("03", "Product & how it works", "The 3-tier suite and data engine"),
    ("04", "Why Sparta leads", "The USP that makes it defensible"),
    ("05", "PetroQuant's challenges", "The real gaps to close"),
    ("06", "The roadmap", "How PetroQuant becomes Sparta-class"),
]
x0, y0, cw, ch, gapx, gapy = 0.75, 2.0, 5.75, 1.35, 0.35, 0.28
for i, (n, h, d) in enumerate(ag):
    col = i % 2; row = i // 2
    x = x0 + col * (cw + gapx); y = y0 + row * (ch + gapy)
    rect(s, x, y, cw, ch, fill=LIGHT, rounded=True)
    rect(s, x, y, 0.09, ch, fill=ACCENT if col == 0 else PETROL2)
    txt(s, x + 0.3, y + 0.22, 1.2, 1.0,
        [[(n, 30, ACCENT if col == 0 else PETROL2, True, FONT_L)]], space_after=0)
    txt(s, x + 1.35, y + 0.24, cw - 1.6, 1.0,
        [[(h, 17, PETROL, True, FONT)],
         [(d, 12.5, MUTED, False, FONT)]], space_after=3)

# =========================================================================
# SLIDE 3 — WHAT IS SPARTA
# =========================================================================
s = slide()
header(s, "01  ·  The company", "What is Sparta Commodities?", 3)
# left: definition card
rect(s, 0.75, 2.0, 6.3, 4.5, fill=PETROL, rounded=True, shadow=True)
txt(s, 1.1, 2.35, 5.7, 0.5,
    [[("The operating system for oil trading desks", 19, WHITE, True, FONT)]],
    space_after=0)
rect(s, 1.1, 3.05, 1.3, 0.04, fill=ACCENT)
txt(s, 1.1, 3.25, 5.75, 3.1,
    [[("A ", 15, RGBColor(0xD6,0xE2,0xEC), False, FONT),
      ("forward-looking market-intelligence & analytics platform",
       15, WHITE, True, FONT),
      (" built specifically for oil traders.", 15, RGBColor(0xD6,0xE2,0xEC), False, FONT)],
     [("", 6, WHITE, False, FONT)],
     [("It connects the ", 15, RGBColor(0xD6,0xE2,0xEC), False, FONT),
      ("physical and paper oil markets", 15, WHITE, True, FONT),
      (" in one live workspace \u2014 pricing, forward curves, arbitrage, freight and blend margins \u2014 so a whole desk trades off one shared, real-time view instead of seven spreadsheets.",
       15, RGBColor(0xD6,0xE2,0xEC), False, FONT)]],
    space_after=6, line_spacing=1.12)
# right: quick identity
items = [
    "Founded 2020|  \u2014 Geneva, Switzerland",
    "By ex-traders|  \u2014 Schuurman (CEO) & Moseley (COO), 37 yrs combined",
    "Category|  \u2014 SaaS market intelligence for commodity desks",
    "The pitch|  \u2014 \u201csee the trade before the market does\u201d",
    "Users|  \u2014 traders, analysts, risk & data teams",
]
rect(s, 7.35, 2.0, 5.2, 4.5, fill=LIGHT, rounded=True)
txt(s, 7.7, 2.3, 4.6, 0.4, [[("IN ONE LINE", 12, ACCENT, True, FONT_SB)]], space_after=0)
bullets(s, 7.7, 2.95, 4.6, 3.4, items, size=14.5, gap=15)

# =========================================================================
# SLIDE 4 — AT A GLANCE (STAT TILES)
# =========================================================================
s = slide()
header(s, "02  ·  Key facts", "Sparta at a glance", 4)
tiles = [
    ("2020", "Founded", ACCENT),
    ("Geneva", "HQ  ·  8 global offices", PETROL2),
    ("$65.5M", "Total funding raised", TEAL),
    ("~200", "Employees  ·  9 countries", PETROL2),
    ("30+", "Broker & data partners", ACCENT),
    ("4 majors", "Trafigura · Chevron · P66 · Equinor", TEAL),
]
x0, y0, cw, ch, gx, gy = 0.75, 2.05, 3.75, 1.95, 0.28, 0.3
for i, (big, lab, c) in enumerate(tiles):
    col = i % 3; row = i // 3
    x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
    rect(s, x, y, cw, ch, fill=CARD, line=LINE, line_w=1, rounded=True, shadow=True)
    rect(s, x, y, cw, 0.10, fill=c)
    txt(s, x + 0.3, y + 0.35, cw - 0.6, 0.9,
        [[(big, 34, PETROL, True, FONT)]], space_after=0)
    txt(s, x + 0.32, y + 1.28, cw - 0.6, 0.6,
        [[(lab, 13, MUTED, False, FONT)]], space_after=0)
txt(s, 0.75, 6.55, 11.8, 0.4,
    [[("Revenue is undisclosed; third-party estimates put ARR at ~$15\u201325M (treat as estimate, not fact).",
       11, MUTED, False, FONT, True)]], space_after=0)

# =========================================================================
# SLIDE 5 — FUNDING JOURNEY
# =========================================================================
s = slide()
header(s, "02  ·  Trajectory", "A five-year funding climb", 5)
steps = [
    ("2020", "Founded", "Two ex-oil traders, Geneva", PETROL2),
    ("Mar 2022", "$6M Seed", "Led by Singular", TEAL),
    ("Oct 2023", "$17.5M Series A", "Led by FirstMark", ACCENT),
    ("Feb 2025", "$42M Series B", "Led by One Peak", PETROL),
]
# timeline base line
rect(s, 1.1, 3.35, 11.1, 0.05, fill=LINE)
x0, cw, gx = 0.9, 2.7, 0.35
for i, (dt, amt, who, c) in enumerate(steps):
    x = x0 + i * (cw + gx)
    # node
    rect(s, x + cw/2 - 0.11, 3.24, 0.22, 0.22, fill=c, rounded=True)
    # card
    rect(s, x, 3.75, cw, 1.85, fill=CARD, line=LINE, line_w=1, rounded=True, shadow=True)
    rect(s, x, 3.75, cw, 0.09, fill=c)
    txt(s, x + 0.25, 4.0, cw - 0.5, 0.4, [[(dt, 12.5, MUTED, True, FONT_SB)]], space_after=0)
    txt(s, x + 0.25, 4.35, cw - 0.5, 0.6, [[(amt, 21, PETROL, True, FONT)]], space_after=0)
    txt(s, x + 0.25, 5.0, cw - 0.5, 0.5, [[(who, 12.5, MUTED, False, FONT)]], space_after=0)
    txt(s, x, 2.5, cw, 0.5, [[(dt.split()[-1], 13, c, True, FONT)]],
        align=PP_ALIGN.CENTER, space_after=0)
txt(s, 0.75, 6.35, 11.8, 0.5,
    [[("Signal:  ", 13, ACCENT, True, FONT_SB),
      ("a growth-equity fund (One Peak, $2B AUM) leading Series B = Sparta had crossed into proven, scaling revenue.",
       13, INK, False, FONT)]], space_after=0)

# =========================================================================
# SLIDE 6 — PRODUCT (3 TIERS)
# =========================================================================
s = slide()
header(s, "03  ·  The product", "One platform, three levels of confidence", 6)
cols = [
    ("Sparta Curves", "\u201cWhere's the market now?\u201d", PETROL2,
     ["Live + historical forward curves",
      "Futures, swaps, physical, freight",
      "Excel plug-in + mobile + sharing"]),
    ("Sparta Knowledge", "\u201cWhy did it move?\u201d", TEAL,
     ["Everything in Curves, plus:",
      "Curated news + expert commentary",
      "Real-time trader calls & context"]),
    ("Sparta Intelligence", "\u201cWhat should I do?\u201d", ACCENT,
     ["Everything above, plus:",
      "AI signals, arbs, blend margins",
      "Scenario testing \u2014 the co-pilot"]),
]
x0, y0, cw, ch, gx = 0.75, 2.05, 3.83, 4.35, 0.3
for i, (name, q, c, its) in enumerate(cols):
    x = x0 + i * (cw + gx)
    rect(s, x, y0, cw, ch, fill=CARD, line=LINE, line_w=1, rounded=True, shadow=True)
    rect(s, x, y0, cw, 1.15, fill=c, rounded=True)
    rect(s, x, y0 + 0.6, cw, 0.55, fill=c)  # square off bottom of header
    txt(s, x + 0.3, y0 + 0.22, cw - 0.6, 0.5, [[(name, 17, WHITE, True, FONT)]], space_after=0)
    txt(s, x + 0.3, y0 + 0.66, cw - 0.6, 0.4, [[(q, 13, WHITE, False, FONT, True)]], space_after=0)
    bullets(s, x + 0.32, y0 + 1.45, cw - 0.6, 2.7, its, size=13.5, gap=13, marker=c)
txt(s, 0.75, 6.62, 11.8, 0.4,
    [[("Modular:  ", 11.5, ACCENT, True, FONT_SB),
      ("each tier works alone or stacks. Add-ons: Data Marketplace (API/SFTP) + S&P Platts integration.",
       11.5, MUTED, False, FONT)]], space_after=0)

# =========================================================================
# SLIDE 7 — HOW IT WORKS
# =========================================================================
s = slide()
header(s, "03  ·  Under the hood", "How Sparta actually does it", 7)
blocks = [
    ("Data engine", ACCENT,
     ["30+ brokerages & data houses, 5+ yr partnerships",
      "Every number verified, auditable, time-aligned"]),
    ("Forward-looking", PETROL2,
     ["Sells opportunities you can still trade",
      "Not backward-looking history like rivals"]),
    ("AI / ML layer", TEAL,
     ["Over/undervaluation scores on 80+ grades & routes",
      "Live arb, freight & blend-margin forecasts"]),
    ("Delivery", PETROL2,
     ["Web dashboard + Excel plug-in + mobile + API",
      "Same live view everywhere the desk works"]),
    ("Human edge", ACCENT,
     ["Insights shaped by ex-trading-desk analysts",
      "Cited by Bloomberg, Reuters & the FT"]),
    ("Collaboration", TEAL,
     ["One shared source of truth for the whole desk",
      "Turns isolated calls into \u201ccollective alpha\u201d"]),
]
x0, y0, cw, ch, gx, gy = 0.75, 2.05, 3.75, 2.05, 0.28, 0.25
for i, (h, c, its) in enumerate(blocks):
    col = i % 3; row = i // 3
    x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
    rect(s, x, y, cw, ch, fill=LIGHT, rounded=True)
    rect(s, x, y, 0.09, ch, fill=c)
    txt(s, x + 0.3, y + 0.22, cw - 0.5, 0.4, [[(h, 15, PETROL, True, FONT)]], space_after=0)
    bullets(s, x + 0.32, y + 0.72, cw - 0.55, 1.2, its, size=11.8, gap=6, marker=c)

# =========================================================================
# SLIDE 8 — WHY SPARTA LEADS (USP)
# =========================================================================
s = slide()
header(s, "04  ·  The moat", "Why Sparta leads the category", 8)
usps = [
    ("Forward-looking & tradable", "The only platform selling the next move, not the last one."),
    ("Physical + paper, unified", "Mirrors how real oil desks actually operate."),
    ("Built by traders, for traders", "Domain credibility tech-only rivals can't fake."),
    ("Collaboration = collective alpha", "One shared truth; sticky team workflow."),
]
x0, y0, cw, ch, gx, gy = 0.75, 2.05, 5.75, 1.7, 0.35, 0.3
for i, (h, d) in enumerate(usps):
    col = i % 2; row = i // 2
    x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
    rect(s, x, y, cw, ch, fill=CARD, line=LINE, line_w=1, rounded=True, shadow=True)
    rect(s, x + 0.28, y + 0.32, 0.55, 0.55, fill=PETROL, rounded=True)
    txt(s, x + 0.28, y + 0.34, 0.55, 0.55, [[(str(i+1), 22, WHITE, True, FONT)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    txt(s, x + 1.05, y + 0.3, cw - 1.3, 1.2,
        [[(h, 16.5, PETROL, True, FONT)],
         [(d, 13, MUTED, False, FONT)]], space_after=4, line_spacing=1.05)
rect(s, 0.75, 5.9, 11.8, 0.85, fill=PETROL, rounded=True)
txt(s, 1.1, 5.9, 11.2, 0.85,
    [[("Result:  ", 14, ACCENT, True, FONT_SB),
      ("blue-chip clients + high switching costs \u2192 strong retention. Its one gap: no public signal-accuracy track record.",
       14, WHITE, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)

# =========================================================================
# SLIDE 9 — PETROQUANT CHALLENGES
# =========================================================================
s = slide()
header(s, "05  ·  The gap", "What PetroQuant must solve \u2014 and why", 9)
rows = [
    ("Data depth", "Public / free feeds only", "Sparta: 30+ exclusive broker partnerships"),
    ("Physical visibility", "No tanker / satellite flow data", "Desks pay for real cargo intelligence"),
    ("Market breadth", "Single market (WTI)", "Sparta spans grades, products & freight"),
    ("Proven edge", "No clean multi-month live proof", "~52% accuracy; edge still fragile"),
    ("Risk design", "Stops erased ~70% of gross profit", "Tail-risk control is the #1 P&L lever"),
    ("Product layer", "Script/dashboard, no delivery", "Sparta: Excel + mobile + API + collab"),
]
# table header
tx, ty, w1, w2, w3 = 0.75, 2.05, 3.0, 4.3, 4.5
rect(s, tx, ty, w1 + w2 + w3, 0.55, fill=PETROL)
for lbl, xx, ww in [("Challenge", tx, w1), ("PetroQuant today", tx + w1, w2),
                    ("Why it matters", tx + w1 + w2, w3)]:
    txt(s, xx + 0.25, ty + 0.11, ww - 0.3, 0.4, [[(lbl, 13, WHITE, True, FONT_SB)]], space_after=0)
for i, (a, b, c) in enumerate(rows):
    y = ty + 0.55 + i * 0.68
    bg = LIGHT if i % 2 == 0 else CARD
    rect(s, tx, y, w1 + w2 + w3, 0.68, fill=bg)
    txt(s, tx + 0.25, y + 0.16, w1 - 0.3, 0.5, [[(a, 13.5, PETROL, True, FONT)]], space_after=0)
    txt(s, tx + w1 + 0.25, y + 0.17, w2 - 0.35, 0.5, [[(b, 12.5, INK, False, FONT)]], space_after=0)
    txt(s, tx + w1 + w2 + 0.25, y + 0.17, w3 - 0.35, 0.5, [[(c, 12.5, MUTED, False, FONT)]], space_after=0)
rect(s, tx, ty, w1 + w2 + w3, 0.55 + len(rows) * 0.68, line=LINE, line_w=1)

# =========================================================================
# SLIDE 10 — ROADMAP
# =========================================================================
s = slide()
header(s, "06  ·  The path", "How PetroQuant becomes Sparta-class", 10)
phases = [
    ("P1", "Prove the edge", "Fix risk & stops · 3-month clean paper · public track record", ACCENT),
    ("P2", "Productize", "Regime dashboard · signal API · Excel & mobile delivery", PETROL2),
    ("P3", "Differentiate", "Win on TRANSPARENCY \u2014 public accuracy scorecard (Sparta's gap)", TEAL),
    ("P4", "Expand", "Add Brent, gold & MCX bridge · richer macro data", PETROL2),
    ("P5", "Community", "Serve the indie / prop / small-fund macro-regime niche", ACCENT),
]
y0 = 2.1
for i, (p, h, d, c) in enumerate(phases):
    y = y0 + i * 0.92
    rect(s, 0.75, y, 0.85, 0.75, fill=c, rounded=True)
    txt(s, 0.75, y + 0.05, 0.85, 0.65, [[(p, 24, WHITE, True, FONT)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    rect(s, 1.75, y, 10.8, 0.75, fill=LIGHT, rounded=True)
    rect(s, 1.75, y, 0.08, 0.75, fill=c)
    txt(s, 2.05, y + 0.12, 3.1, 0.55, [[(h, 16, PETROL, True, FONT)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    txt(s, 5.1, y + 0.13, 7.2, 0.55, [[(d, 12.8, MUTED, False, FONT)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=0)

# =========================================================================
# SLIDE 11 — CLOSING
# =========================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=PETROL)
rect(s, 0, 0, 13.333, 0.14, fill=ACCENT)
rect(s, 0, 7.36, 13.333, 0.14, fill=ACCENT)
txt(s, 0.9, 1.5, 11.5, 0.5, [[("THE TAKEAWAY", 13, ACCENT, True, FONT_SB)]], space_after=0)
txt(s, 0.85, 2.15, 11.7, 2.2,
    [[("Don't out-gun Kpler on data.", 34, WHITE, True, FONT)],
     [("Out-trust everyone on transparency.", 34, WHITE, True, FONT)]],
    space_after=6, line_spacing=1.05)
txt(s, 0.9, 4.2, 11.2, 1.2,
    [[("Sparta owns the physical-desk workflow. PetroQuant's opening is the ",
       16, RGBColor(0xC7,0xD6,0xE2), False, FONT),
      ("verifiable, forward-looking co-pilot", 16, WHITE, True, FONT),
      (" for oil-macro traders \u2014 the niche where an auditable track record beats scale.",
       16, RGBColor(0xC7,0xD6,0xE2), False, FONT)]],
    space_after=0, line_spacing=1.15)
rect(s, 0.9, 5.75, 3.4, 0.05, fill=ACCENT)
txt(s, 0.9, 5.95, 11, 0.5,
    [[("PetroQuant  ·  Competitive Strategy Brief  ·  July 2026",
       12, RGBColor(0xAE,0xC0,0xCE), False, FONT)]], space_after=0)

prs.save("PetroQuant_Sparta_Strategy.pptx")
print("SAVED PetroQuant_Sparta_Strategy.pptx  |  slides:", len(prs.slides._sldIdLst))
